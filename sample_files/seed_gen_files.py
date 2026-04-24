import os
import json
import csv
import zipfile
import shutil
from pathlib import Path
from datetime import datetime, timedelta

# External libs
from docx import Document
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

# -----------------------------
# CONFIG
# -----------------------------
OUTPUT_DIR = Path("langchain_large_learning_corpus")
ZIP_NAME = "langchain_large_learning_corpus.zip"

# Large file sizes
LARGE_TEXT_MB = [5, 10, 20]
LARGE_JSONL_MB = [5, 10]
LARGE_CSV_MB = [5, 15]
LARGE_PDF_PAGES = [100, 200]
LARGE_DOCX_PARAGRAPHS = [800, 1500]
LARGE_XLSX_ROWS = [10000, 25000]
LARGE_PPTX_SLIDES = [30, 60]

# -----------------------------
# HELPERS
# -----------------------------
def reset_dir(path: Path):
    if path.exists():
        shutil.rmtree(path)
    path.mkdir(parents=True, exist_ok=True)

def write_text(path: Path, content: str):
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")

def repeat_paragraph(topic: str, idx: int) -> str:
    return (
        f"{topic} paragraph {idx}. "
        f"This file is generated for LangChain learning, document loading, parsing, "
        f"text splitting, embedding creation, vector indexing, retrieval, metadata extraction, "
        f"semantic search, summarization, and RAG experiments. "
        f"It contains structured repeated content so that the file becomes large enough for testing "
        f"while still remaining meaningful for chunking and retrieval workflows. "
        f"Each paragraph includes consistent language patterns that are useful for evaluating "
        f"document loaders, overlap behavior, filtering, and search relevance.\n"
    )

def generate_large_text(target_mb: int, topic: str) -> str:
    target_bytes = target_mb * 1024 * 1024
    parts = []
    current = 0
    i = 1
    while current < target_bytes:
        para = repeat_paragraph(topic, i)
        parts.append(para)
        current += len(para.encode("utf-8"))
        i += 1
    return "".join(parts)

def write_large_txt(path: Path, target_mb: int, topic: str):
    content = generate_large_text(target_mb, topic)
    write_text(path, content)

def write_large_jsonl(path: Path, target_mb: int, topic: str):
    path.parent.mkdir(parents=True, exist_ok=True)
    target_bytes = target_mb * 1024 * 1024
    current = 0
    i = 1
    with open(path, "w", encoding="utf-8") as f:
        while current < target_bytes:
            obj = {
                "id": i,
                "topic": topic,
                "title": f"{topic} record {i}",
                "summary": f"Synthetic record {i} for large-scale LangChain testing.",
                "content": repeat_paragraph(topic, i) * 3,
                "tags": ["langchain", "rag", "retrieval", "chunking", f"batch-{i % 10}"],
                "owner": f"user_{i % 25}",
                "priority": ["low", "medium", "high"][i % 3],
                "active": i % 2 == 0,
            }
            line = json.dumps(obj, ensure_ascii=False) + "\n"
            f.write(line)
            current += len(line.encode("utf-8"))
            i += 1

def write_large_csv(path: Path, target_mb: int, topic: str):
    path.parent.mkdir(parents=True, exist_ok=True)
    target_bytes = target_mb * 1024 * 1024
    current = 0
    i = 1
    header = ["id", "name", "category", "description", "status", "owner", "region", "score"]
    with open(path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(header)
        current += len((",".join(header) + "\n").encode("utf-8"))
        while current < target_bytes:
            row = [
                i,
                f"{topic} Item {i}",
                f"Category-{i % 12}",
                repeat_paragraph(topic, i).replace("\n", " ").strip(),
                "active" if i % 2 == 0 else "inactive",
                f"user_{i % 50}",
                f"region_{i % 5}",
                round((i % 100) / 10, 2),
            ]
            line = ",".join(map(lambda x: str(x).replace(",", " "), row)) + "\n"
            f.write(line)
            current += len(line.encode("utf-8"))
            i += 1

def write_large_pdf(path: Path, pages: int, topic: str):
    path.parent.mkdir(parents=True, exist_ok=True)
    c = canvas.Canvas(str(path), pagesize=A4)
    w, h = A4

    para_index = 1
    for p in range(1, pages + 1):
        c.setFont("Helvetica-Bold", 14)
        c.drawString(50, h - 50, f"{topic} - Page {p}")

        c.setFont("Helvetica", 10)
        y = h - 80
        while y > 50:
            line = repeat_paragraph(topic, para_index)[:130]
            c.drawString(50, y, line)
            y -= 14
            para_index += 1

        c.showPage()

    c.save()

def write_large_docx(path: Path, paragraphs: int, topic: str):
    path.parent.mkdir(parents=True, exist_ok=True)
    doc = Document()
    doc.add_heading(topic, 0)
    doc.add_paragraph("Large DOCX generated for LangChain loader and parser testing.")
    for i in range(1, paragraphs + 1):
        doc.add_paragraph(repeat_paragraph(topic, i) * 2)
    doc.save(path)

def write_large_xlsx(path: Path, rows: int, topic: str):
    path.parent.mkdir(parents=True, exist_ok=True)
    wb = Workbook()
    ws = wb.active
    ws.title = "LargeData"
    ws.append(["id", "title", "description", "owner", "status", "score"])
    for i in range(1, rows + 1):
        ws.append([
            i,
            f"{topic} Record {i}",
            repeat_paragraph(topic, i).replace("\n", " ")[:500],
            f"user_{i % 100}",
            "open" if i % 2 else "closed",
            i % 100
        ])
    wb.save(path)

def write_large_pptx(path: Path, slides_count: int, topic: str):
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    for i in range(1, slides_count + 1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"{topic} Slide {i}"
        slide.placeholders[1].text = (
            f"{repeat_paragraph(topic, i)[:300]}\n\n"
            f"Highlights:\n"
            f"- LangChain testing\n"
            f"- Chunking and retrieval\n"
            f"- Metadata extraction\n"
            f"- Synthetic learning content"
        )
    prs.save(path)

def create_zip(source_dir: Path, zip_path: Path):
    if zip_path.exists():
        zip_path.unlink()
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for file_path in sorted(source_dir.rglob("*")):
            if file_path.is_file():
                zf.write(file_path, arcname=str(file_path.relative_to(source_dir)))

# -----------------------------
# GENERATORS
# -----------------------------
def create_readme(base: Path):
    content = f"""# LangChain Large Learning Corpus

Generated on: {datetime.now().isoformat(timespec="seconds")}

This dataset contains 100+ files across many formats, including both:
- small realistic samples
- large files for stress testing

Useful for:
- document loaders
- parsers
- chunking
- embeddings
- vector stores
- retrieval
- filtering
- metadata extraction
- RAG experiments

Main categories:
- text
- markdown
- data
- code
- config
- markup
- logs
- emails
- office files
- large files

"""
    write_text(base / "README.md", content)

def create_small_text_files(base: Path):
    folder = base / "text"
    samples = {
        "customer_support_notes.txt": """Customer Support Notes

Issue categories:
1. Login failures after password reset
2. PDF export timing out for large reports
3. Search returning stale cache results
4. Users asking for dark mode
""",
        "meeting_transcript.txt": """Meeting Transcript

Asha: We need clearer ingestion for mixed file formats.
Rahul: Let's separate parsing, cleaning, and chunking.
Asha: Also capture source metadata.
""",
        "incident_report.txt": """Incident Report

Date: 2026-04-23
Severity: Medium
Service: document-indexer
Summary: Some XML files failed to parse due to malformed encoding.
""",
        "faq_plain_text.txt": """FAQ

Q: What is chunk overlap?
A: It preserves context between adjacent chunks.

Q: Why store metadata?
A: It helps filtering and citation.
""",
        "product_requirements.txt": """Product Requirements Summary

Project: Knowledge Assistant
Core features:
- ingest multiple document types
- semantic retrieval
- citations
- filtering
""",
    }
    for name, content in samples.items():
        write_text(folder / name, content)

    for i in range(1, 21):
        write_text(
            folder / f"note_{i:02d}.txt",
            f"Sample text note {i}\n\n{repeat_paragraph('LangChain note', i) * 3}"
        )

def create_markdown_files(base: Path):
    folder = base / "markdown"
    for i in range(1, 21):
        content = f"""# Lesson {i}

## Topic
LangChain concept {i}

## Explanation
{repeat_paragraph('Markdown lesson', i)}

## Use Case
- loaders
- splitters
- retrieval
- RAG
"""
        write_text(folder / f"lesson_{i:02d}.md", content)

def create_data_files(base: Path):
    folder = base / "data"

    # CSV
    with open(folder / "products.csv", "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["product_id", "name", "category", "price", "stock"])
        for i in range(1, 101):
            writer.writerow([f"P{i:03d}", f"Product {i}", f"Category {i % 10}", i * 100, i * 5])

    # TSV
    with open(folder / "support_tickets.tsv", "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f, delimiter="\t")
        writer.writerow(["ticket_id", "priority", "status", "summary"])
        for i in range(1, 101):
            writer.writerow([f"T{i:04d}", ["Low", "Medium", "High"][i % 3], "Open", f"Ticket summary {i}"])

    # JSON / JSONL / NDJSON / YAML / TOML / XML
    write_text(folder / "inventory.json", json.dumps([
        {"sku": f"SKU-{i:03d}", "name": f"Item {i}", "qty": i * 2}
        for i in range(1, 51)
    ], indent=2))

    write_text(folder / "users.jsonl", "\n".join(
        json.dumps({"id": i, "name": f"user_{i}", "role": ["admin", "viewer", "editor"][i % 3]})
        for i in range(1, 101)
    ))

    write_text(folder / "events.ndjson", "\n".join(
        json.dumps({"event_id": i, "event": "query", "user": f"user_{i % 10}"})
        for i in range(1, 101)
    ))

    write_text(folder / "countries.yaml", """countries:
  - code: IN
    name: India
    region: Asia
  - code: US
    name: United States
    region: North America
""")

    write_text(folder / "sales_report.toml", """[summary]
quarter = "Q1-2026"
revenue = 1250000
currency = "INR"
""")

    write_text(folder / "metrics.xml", """<?xml version="1.0" encoding="UTF-8"?>
<metrics>
  <service name="ingest">
    <latency_p95_ms>840</latency_p95_ms>
  </service>
</metrics>
""")

    for i in range(1, 16):
        write_text(folder / f"sample_data_{i:02d}.json", json.dumps({
            "id": i,
            "title": f"Sample record {i}",
            "content": repeat_paragraph("JSON sample", i),
            "score": i / 10
        }, indent=2))

def create_config_files(base: Path):
    folder = base / "config"
    config_files = {
        ".env": "OPENAI_API_KEY=demo-key\nAPP_ENV=development\nTOP_K=5\n",
        "settings.ini": "[app]\nname=rag_demo\nmode=local\n",
        "application.properties": "server.port=8000\napp.name=knowledge-assistant\n",
        "pipeline.yaml": "loader: unstructured\nsplitter: recursive\nvector_store: chroma\n",
        "service.toml": 'name="query-service"\nversion="0.1.0"\n',
        "parser_rules.conf": "allow_pdf=true\nallow_docx=true\nmax_file_size_mb=100\n",
        "feature_flags.yml": "flags:\n  enable_reranker: true\n  enable_citations: true\n",
        "mappings.json": json.dumps({"txt": "text/plain", "md": "text/markdown"}, indent=2),
    }
    for name, content in config_files.items():
        write_text(folder / name, content)

    for i in range(1, 16):
        write_text(folder / f"profile_{i:02d}.env", f"PROFILE=profile_{i:02d}\nCACHE_TTL={i*10}\n")

def create_markup_files(base: Path):
    folder = base / "markup"

    write_text(folder / "landing_page.html", """<!doctype html>
<html><head><title>Knowledge Portal</title></head>
<body><h1>Knowledge Portal</h1><p>Portal for policies and notes.</p></body></html>""")

    write_text(folder / "catalog.xml", """<?xml version="1.0"?>
<catalog><book id="b1"><title>IR Basics</title></book></catalog>""")

    write_text(folder / "diagram.svg", """<svg xmlns="http://www.w3.org/2000/svg" width="300" height="100">
<rect x="10" y="10" width="80" height="30" fill="#ddd" stroke="#333"/>
<text x="25" y="30">Loader</text>
</svg>""")

    write_text(folder / "paper.tex", r"""\documentclass{article}
\begin{document}
\section*{Sample LaTeX}
This is a sample file.
\end{document}
""")

    write_text(folder / "guide.rst", """Guide
=====

Overview
--------

Sample RST document.
""")

    for i in range(1, 16):
        write_text(folder / f"web_fragment_{i:02d}.html",
                   f"<html><body><h3>Fragment {i}</h3><p>{repeat_paragraph('HTML fragment', i)}</p></body></html>")

def create_log_files(base: Path):
    folder = base / "logs"
    start = datetime(2026, 4, 23, 9, 0, 0)

    for i in range(1, 16):
        lines = []
        t = start + timedelta(minutes=i)
        for j in range(1, 51):
            ts = (t + timedelta(seconds=j * 10)).isoformat()
            level = ["INFO", "WARN", "ERROR"][j % 3]
            lines.append(f"{ts} {level} service=document-indexer request_id=req-{i:02d}-{j:03d} event=sample_event_{j}")
        write_text(folder / f"app_log_{i:02d}.log", "\n".join(lines))

def create_email_files(base: Path):
    folder = base / "emails"
    for i in range(1, 16):
        eml = f"""From: sender{i}@example.com
To: learner@example.com
Subject: Sample Email {i}
Date: Wed, 23 Apr 2026 10:{i:02d}:00 +0530
MIME-Version: 1.0
Content-Type: text/plain; charset="UTF-8"

Hello,
This is sample email number {i} for LangChain learning.
Regards,
Demo Team
"""
        write_text(folder / f"sample_email_{i:02d}.eml", eml)

def create_code_files(base: Path):
    code_samples = {
        ("code/python", "app.py"): """from fastapi import FastAPI
app = FastAPI()

@app.get("/health")
def health():
    return {"status": "ok"}
""",
        ("code/python", "chunking_demo.py"): """def split_text(text, size=200):
    return [text[i:i+size] for i in range(0, len(text), size)]
""",
        ("code/javascript", "server.js"): """const http = require('http');
http.createServer((req, res) => { res.end('ok'); }).listen(3000);
""",
        ("code/java", "Main.java"): """public class Main {
  public static void main(String[] args) {
    System.out.println("Hello LangChain");
  }
}
""",
        ("code/c_cpp", "main.c"): """#include <stdio.h>
int main() { printf("hello\\n"); return 0; }
""",
        ("code/c_cpp", "sort.cpp"): """#include <vector>
#include <algorithm>
int main() { std::vector<int> v={3,1,2}; std::sort(v.begin(), v.end()); }
""",
        ("code/misc", "script.sh"): "#!/bin/bash\necho \"ingest started\"\n",
        ("code/misc", "build.bat"): "@echo off\necho Building project\n",
        ("code/misc", "deploy.ps1"): "Write-Output \"Deploying demo stack\"\n",
        ("code/misc", "app.go"): "package main\nimport \"fmt\"\nfunc main(){ fmt.Println(\"hello\") }\n",
        ("code/misc", "analyze.rb"): "puts 'ruby sample'\n",
        ("code/misc", "index.php"): "<?php echo 'php sample'; ?>\n",
        ("code/misc", "service.kt"): "fun main(){ println(\"kotlin sample\") }\n",
        ("code/misc", "ios.swift"): "print(\"swift sample\")\n",
        ("code/misc", "analysis.r"): "x <- c(1,2,3)\nmean(x)\n",
        ("code/misc", "query.sql"): "SELECT * FROM tickets WHERE priority = 'High';\n",
        ("code/misc", "style.css"): "body { font-family: Arial; margin: 16px; }\n",
    }

    for (folder, name), content in code_samples.items():
        write_text(base / folder / name, content)

    for i in range(1, 21):
        write_text(base / "code/python" / f"module_{i:02d}.py", f"def func_{i}():\n    return 'module {i}'\n")
        write_text(base / "code/javascript" / f"widget_{i:02d}.ts", f"export const item{i}: string = 'widget {i}';\n")

def create_docs_files(base: Path):
    folder = base / "docs"
    for i in range(1, 21):
        content = f"""# Case Study {i}

Problem:
Improve search quality for dataset {i}.

Approach:
Use better chunking, metadata enrichment, and reranking.

Outcome:
Synthetic improvement in retrieval quality.
"""
        write_text(folder / f"case_study_{i:02d}.md", content)

def create_media_like_files(base: Path):
    folder = base / "media_like"
    write_text(folder / "video_subtitles.srt", """1
00:00:00,000 --> 00:00:02,000
Welcome to the retrieval tutorial.

2
00:00:02,100 --> 00:00:05,000
In this lesson, we compare chunk sizes and overlaps.
""")

    write_text(folder / "meeting_captions.vtt", """WEBVTT

00:00:00.000 --> 00:00:02.000
Today we review document ingestion metrics.

00:00:02.100 --> 00:00:05.000
Next, we discuss metadata normalization.
""")

def create_small_office_files(base: Path):
    folder = base / "office"

    # DOCX
    doc = Document()
    doc.add_heading("Sample Project Proposal", 0)
    doc.add_paragraph("This DOCX file is included for loader experiments.")
    doc.add_heading("Objectives", level=1)
    doc.add_paragraph("Ingest mixed formats", style="List Bullet")
    doc.add_paragraph("Store metadata", style="List Bullet")
    doc.add_paragraph("Enable retrieval with citations", style="List Bullet")
    doc.save(folder / "sample_proposal.docx")

    # XLSX
    wb = Workbook()
    ws = wb.active
    ws.title = "Sales"
    ws.append(["Month", "Revenue", "Cost", "Profit"])
    ws.append(["Jan", 120000, 80000, 40000])
    ws.append(["Feb", 140000, 90000, 50000])
    wb.save(folder / "quarterly_sales.xlsx")

    # PPTX
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Knowledge Assistant Demo"
    slide.placeholders[1].text = "Mixed-format ingestion\nSemantic retrieval\nSource citations"
    prs.save(folder / "demo_deck.pptx")

    # PDF
    pdf_path = folder / "sample_report.pdf"
    c = canvas.Canvas(str(pdf_path), pagesize=A4)
    w, h = A4
    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, h - 72, "Sample PDF Report")
    c.setFont("Helvetica", 11)
    y = h - 110
    for line in [
        "This PDF is a small sample for LangChain learning.",
        "It can be used for PDF loader testing.",
        "It contains simple plain text lines."
    ]:
        c.drawString(72, y, line)
        y -= 18
    c.showPage()
    c.save()

def create_large_files(base: Path):
    folder = base / "large_files"

    # Large TXT
    for size in LARGE_TEXT_MB:
        write_large_txt(folder / f"large_text_{size}mb.txt", size, f"Large text {size}MB")

    # Large JSONL
    for size in LARGE_JSONL_MB:
        write_large_jsonl(folder / f"large_jsonl_{size}mb.jsonl", size, f"Large JSONL {size}MB")

    # Large CSV
    for size in LARGE_CSV_MB:
        write_large_csv(folder / f"large_csv_{size}mb.csv", size, f"Large CSV {size}MB")

    # Large PDF
    for pages in LARGE_PDF_PAGES:
        write_large_pdf(folder / f"large_pdf_{pages}_pages.pdf", pages, f"Large PDF {pages} pages")

    # Large DOCX
    for paras in LARGE_DOCX_PARAGRAPHS:
        write_large_docx(folder / f"large_docx_{paras}_paras.docx", paras, f"Large DOCX {paras} paras")

    # Large XLSX
    for rows in LARGE_XLSX_ROWS:
        write_large_xlsx(folder / f"large_xlsx_{rows}_rows.xlsx", rows, f"Large XLSX {rows} rows")

    # Large PPTX
    for slides in LARGE_PPTX_SLIDES:
        write_large_pptx(folder / f"large_pptx_{slides}_slides.pptx", slides, f"Large PPTX {slides} slides")

def create_manifest(base: Path):
    files = []
    for p in sorted(base.rglob("*")):
        if p.is_file():
            files.append(str(p.relative_to(base)))
    write_text(base / "manifest.txt", "\n".join(files) + "\n")
    return len(files)

# -----------------------------
# MAIN
# -----------------------------
def main():
    reset_dir(OUTPUT_DIR)

    # Create folders
    folders = [
        "text",
        "markdown",
        "data",
        "config",
        "markup",
        "logs",
        "emails",
        "docs",
        "media_like",
        "office",
        "large_files",
        "code/python",
        "code/javascript",
        "code/java",
        "code/c_cpp",
        "code/misc",
    ]
    for folder in folders:
        (OUTPUT_DIR / folder).mkdir(parents=True, exist_ok=True)

    create_readme(OUTPUT_DIR)
    create_small_text_files(OUTPUT_DIR)
    create_markdown_files(OUTPUT_DIR)
    create_data_files(OUTPUT_DIR)
    create_config_files(OUTPUT_DIR)
    create_markup_files(OUTPUT_DIR)
    create_log_files(OUTPUT_DIR)
    create_email_files(OUTPUT_DIR)
    create_code_files(OUTPUT_DIR)
    create_docs_files(OUTPUT_DIR)
    create_media_like_files(OUTPUT_DIR)
    create_small_office_files(OUTPUT_DIR)
    create_large_files(OUTPUT_DIR)

    total_files = create_manifest(OUTPUT_DIR)
    create_zip(OUTPUT_DIR, Path(ZIP_NAME))

    print(f"Created folder: {OUTPUT_DIR.resolve()}")
    print(f"Created zip: {Path(ZIP_NAME).resolve()}")
    print(f"Total files generated: {total_files}")

if __name__ == "__main__":
    main()